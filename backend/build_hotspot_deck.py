"""Bangun paparan .pptx: "Analisis Kebakaran Hutan dan Lahan pada Areal
Perhutanan Sosial dan Hutan Adat" — rekapitulasi titik panas, kompleks
kebakaran, dan luas areal kebakaran menurut data Kementerian Kehutanan.
Periode 1 Januari 2026 s.d. tanggal berjalan. Register bahasa: kedinasan.

Sumber angka: extract_hotspot_deck_stats.py -> hotspot_deck.json (READ-ONLY).

Jalankan dari backend/:
    .venv/bin/python build_hotspot_deck.py <hotspot_deck.json> <output.pptx>
"""

from __future__ import annotations

import json
import sys
from datetime import date
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palet & tema
PAPER = RGBColor(0xFA, 0xF8, 0xF5)
INK = RGBColor(0x1A, 0x1D, 0x21)
INK_SOFT = RGBColor(0x4A, 0x55, 0x68)
EMBER = RGBColor(0xE8, 0x64, 0x1A)
RUST = RGBColor(0xB5, 0x24, 0x1C)
GOLD = RGBColor(0xE0, 0x99, 0x22)
SLATE = RGBColor(0x63, 0x72, 0x87)
GREEN = RGBColor(0x2E, 0x6F, 0x4E)
LINE = RGBColor(0xE4, 0xDE, 0xD5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEX = lambda c: "#%02x%02x%02x" % (c[0], c[1], c[2])
FONT = "Lato"

for _p in fm.findSystemFonts():
    if "lato" in _p.lower():
        try:
            fm.fontManager.addfont(_p)
        except Exception:
            pass
plt.rcParams.update({
    "font.family": FONT,
    "axes.edgecolor": HEX(LINE),
    "axes.linewidth": 0.8,
    "text.color": HEX(INK),
    "axes.labelcolor": HEX(INK_SOFT),
    "xtick.color": HEX(INK_SOFT),
    "ytick.color": HEX(INK_SOFT),
    "figure.facecolor": HEX(PAPER),
    "axes.facecolor": HEX(PAPER),
    "savefig.facecolor": HEX(PAPER),
})

SCRATCH = Path("/tmp/claude-1000/-home-ryandshinevps-etaseneu/"
               "16e17ccb-8c0f-40ce-ad1f-be6367b1c126/scratchpad")
IMG = SCRATCH / "img_hs"
IMG.mkdir(parents=True, exist_ok=True)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
DOC_FOOT = ("ETA SENEU  ·  Analisis Kebakaran Hutan dan Lahan pada Areal "
            "Perhutanan Sosial dan Hutan Adat  ·  2026")
BULAN = {1: "Jan", 2: "Feb", 3: "Mar", 4: "Apr", 5: "Mei", 6: "Jun",
         7: "Jul", 8: "Agt", 9: "Sep", 10: "Okt", 11: "Nov", 12: "Des"}
BULAN_PANJANG = {1: "Januari", 2: "Februari", 3: "Maret", 4: "April", 5: "Mei",
                 6: "Juni", 7: "Juli", 8: "Agustus", 9: "September", 10: "Oktober",
                 11: "November", 12: "Desember"}


def ribu(n) -> str:
    return f"{float(n):,.0f}".replace(",", ".")


def ha(x) -> str:
    s = f"{float(x):,.1f}"
    return s.replace(",", "§").replace(".", ",").replace("§", ".")


def pct(a, b) -> str:
    return f"{(100.0 * a / b):.1f}%".replace(".", ",") if b else "0%"


def tanggal_id(iso: str) -> str:
    y, m, d = (int(x) for x in iso.split("-"))
    return f"{d} {BULAN_PANJANG[m]} {y}"


# ---------------------------------------------------------------- helper pptx
def _txt(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.12, space_after=Pt(4)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = space_after
        for (text, size, bold, color, *rest) in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            if rest and rest[0]:
                r.font.italic = True
    return tb


def _rect(slide, x, y, w, h, fill, *, line=None, line_w=Pt(0.75), rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if rounded:
        try:
            shp.adjustments[0] = 0.055
        except Exception:
            pass
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    return shp


def _bg(slide, color):
    _rect(slide, 0, 0, EMU_W, EMU_H, color)


def footer(slide, page, *, dark=False):
    c = RGBColor(0x8A, 0x94, 0xA0) if dark else RGBColor(0x9A, 0xA0, 0xA8)
    _txt(slide, Inches(0.75), Inches(7.04), Inches(10.6), Inches(0.36),
         [[(DOC_FOOT, 7.5, False, c)]])
    _txt(slide, Inches(12.2), Inches(7.04), Inches(0.7), Inches(0.36),
         [[(str(page), 8, True, c)]], align=PP_ALIGN.RIGHT)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def divider(prs, kicker, title, sub, page):
    s = blank(prs)
    _bg(s, INK)
    _rect(s, Inches(0.75), Inches(3.05), Inches(0.9), Inches(0.09), EMBER)
    _txt(s, Inches(0.75), Inches(2.35), Inches(11), Inches(0.5),
         [[(kicker.upper(), 13, True, EMBER)]])
    _txt(s, Inches(0.72), Inches(3.35), Inches(11.8), Inches(1.4),
         [[(title, 38, True, WHITE)]], line_spacing=1.05)
    _txt(s, Inches(0.75), Inches(4.95), Inches(11.5), Inches(1.2),
         [[(sub, 13.5, False, RGBColor(0xB8, 0xC0, 0xC8))]], line_spacing=1.25)
    footer(s, page, dark=True)
    return s


def content(prs, kicker, headline, page):
    s = blank(prs)
    _bg(s, PAPER)
    _rect(s, Inches(0.55), Inches(0.52), Inches(0.07), Inches(0.86), EMBER)
    _txt(s, Inches(0.78), Inches(0.46), Inches(11.9), Inches(0.32),
         [[(kicker.upper(), 11, True, EMBER)]])
    _txt(s, Inches(0.76), Inches(0.74), Inches(12.1), Inches(0.9),
         [[(headline, 21, True, INK)]], line_spacing=1.03)
    _rect(s, Inches(0.78), Inches(1.62), Inches(11.9), Pt(1.4), LINE)
    footer(s, page)
    return s


def stat_cards(slide, items, *, y=Inches(1.9), x0=Inches(0.78), w_total=Inches(11.9),
               h=Inches(1.55)):
    gap = Inches(0.2)
    n = len(items)
    w = Emu(int((w_total - gap * (n - 1)) / n))
    for i, (val, label, accent) in enumerate(items):
        x = Emu(int(x0 + i * (w + gap)))
        _rect(slide, x, y, w, h, WHITE, line=LINE, line_w=Pt(1), rounded=True)
        _rect(slide, x, y, Inches(0.07), h, accent)
        vsize = 25 if len(val) <= 6 else (21 if len(val) <= 9 else 17)
        _txt(slide, Emu(int(x + Inches(0.22))), Emu(int(y + Inches(0.16))),
             Emu(int(w - Inches(0.34))), Inches(0.6), [[(val, vsize, True, accent)]])
        _txt(slide, Emu(int(x + Inches(0.22))), Emu(int(y + Inches(0.74))),
             Emu(int(w - Inches(0.34))), Inches(0.78),
             [[(label, 9, False, INK_SOFT)]], line_spacing=1.08)


def bullets(slide, items, *, x=Inches(0.78), y=Inches(3.7), w=Inches(11.9),
            h=Inches(3.0), size=11.5, gap=Pt(7)):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            paras.append([("▪  ", size, True, EMBER), (lead, size, True, INK),
                          (rest, size, False, INK_SOFT)])
        else:
            paras.append([("▪  ", size, True, EMBER), (it, size, False, INK_SOFT)])
    _txt(slide, x, y, w, h, paras, line_spacing=1.16, space_after=gap)


def _num_cols(rows):
    if not rows:
        return 0
    cnt = 0
    for v in rows[0][::-1]:
        sv = str(v).replace(".", "").replace(",", "").replace("—", "").replace(" ", "")
        if isinstance(v, (int, float)) or sv.isdigit():
            cnt += 1
        else:
            break
    return cnt


def table(slide, headers, rows, *, x=Inches(0.78), y=Inches(1.9), w=Inches(11.9),
          row_h=Inches(0.34), col_w=None, zebra=RGBColor(0xF1, 0xEC, 0xE3), size=10):
    nr, nc = len(rows) + 1, len(headers)
    numc = _num_cols(rows)
    gt = slide.shapes.add_table(nr, nc, x, y, w, Emu(int(row_h * nr))).table
    gt.first_row = False
    gt.horz_banding = False
    if col_w:
        for j, cw in enumerate(col_w):
            gt.columns[j].width = cw
    for j, htext in enumerate(headers):
        c = gt.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = INK
        c.margin_left = c.margin_right = Inches(0.08)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if j >= nc - numc else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = str(htext)
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else zebra
            c.margin_left = c.margin_right = Inches(0.08)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j >= nc - numc else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = INK if j == 0 else INK_SOFT
    for i in range(nr):
        gt.rows[i].height = row_h
    return gt


def place(slide, res, x, y, w):
    path, fw, fh = res
    h = Emu(int(w * fh / fw))
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    return h


# ---------------------------------------------------------------- chart builders
def _finish(fig, name):
    for ax in fig.axes:
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
    fig.tight_layout(pad=0.6)
    w, h = fig.get_size_inches()
    p = IMG / name
    fig.savefig(p, dpi=200)
    plt.close(fig)
    return (p, float(w), float(h))


def chart_monthly(months_vals, name, *, ylabel, color=EMBER, hi_month=8,
                  figsize=(11.8, 3.2), fmt=ribu):
    fig, ax = plt.subplots(figsize=figsize)
    xs = [BULAN[m] for m, _ in months_vals]
    ys = [v for _, v in months_vals]
    cols = [HEX(RUST) if m == hi_month else HEX(color) for m, _ in months_vals]
    bars = ax.bar(xs, ys, color=cols, width=0.62)
    ax.set_ylabel(ylabel, fontsize=10)
    ax.grid(axis="y", color=HEX(LINE), linewidth=0.7)
    ax.set_axisbelow(True)
    for b, v in zip(bars, ys):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height(), fmt(v),
                ha="center", va="bottom", fontsize=9, color=HEX(INK), fontweight="bold")
    ax.margins(y=0.22)
    return _finish(fig, name)


def chart_hbar(labels, values, name, *, color=EMBER, xlabel="", figsize=(11.4, 3.6),
               fmt=ribu, values2=None, legend=None, color2=SLATE):
    fig, ax = plt.subplots(figsize=figsize)
    y = list(range(len(labels)))
    if values2 is None:
        ax.barh(y, values, color=HEX(color), height=0.62)
        for i, v in enumerate(values):
            ax.text(v, i, " " + fmt(v), va="center", fontsize=9,
                    color=HEX(INK), fontweight="bold")
    else:
        hh = 0.38
        ax.barh([i + hh / 2 for i in y], values, color=HEX(color), height=hh,
                label=legend[0] if legend else None)
        ax.barh([i - hh / 2 for i in y], values2, color=HEX(color2), height=hh,
                label=legend[1] if legend else None)
        for i, v in enumerate(values):
            ax.text(v, i + hh / 2, " " + fmt(v), va="center", fontsize=8,
                    color=HEX(INK), fontweight="bold")
        for i, v in enumerate(values2):
            ax.text(v, i - hh / 2, " " + fmt(v), va="center", fontsize=8,
                    color=HEX(INK_SOFT))
        if legend:
            ax.legend(frameon=False, fontsize=9.5, loc="lower right")
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=10)
    ax.invert_yaxis()
    ax.set_xlabel(xlabel, fontsize=9.5)
    ax.grid(axis="x", color=HEX(LINE), linewidth=0.7)
    ax.set_axisbelow(True)
    ax.margins(x=0.18)
    return _finish(fig, name)


def chart_hours(hour_vals, name, *, figsize=(11.8, 2.9)):
    fig, ax = plt.subplots(figsize=figsize)
    hv = {int(h): v for h, v in hour_vals}
    pairs = [(h, hv.get(h, 0)) for h in range(24)]
    xs = [f"{h:02d}" for h, _ in pairs]
    ys = [v for _, v in pairs]
    cols = [HEX(EMBER) if 9 <= h <= 15 else (HEX(SLATE) if (h <= 2 or h >= 22)
            else HEX(LINE)) for h, _ in pairs]
    bars = ax.bar(xs, ys, color=cols, width=0.7)
    ax.set_xlabel("jam (WIB)", fontsize=9.5)
    ax.set_ylabel("titik panas", fontsize=9.5)
    ax.grid(axis="y", color=HEX(LINE), linewidth=0.7)
    ax.set_axisbelow(True)
    ax.margins(y=0.28)
    for b, v in zip(bars, ys):
        if v <= 0:
            continue
        ax.text(b.get_x() + b.get_width() / 2, b.get_height(), ribu(v),
                ha="center", va="bottom", fontsize=7, rotation=90,
                color=HEX(INK), fontweight="bold")
    return _finish(fig, name)


def chart_klhk_monthly(kmv, name, *, figsize=(11.8, 3.1)):
    fig, ax = plt.subplots(figsize=figsize)
    xs = [BULAN[m] for m, _ in kmv] + ["Agt", "Sep"]
    ys = [v for _, v in kmv] + [0, 0]
    cols = [HEX(SLATE)] * len(kmv) + [HEX(LINE), HEX(LINE)]
    bars = ax.bar(xs, ys, color=cols, width=0.62)
    for b, v in zip(bars[:len(kmv)], ys[:len(kmv)]):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height(), ha(v),
                ha="center", va="bottom", fontsize=9, color=HEX(INK), fontweight="bold")
    for b in bars[len(kmv):]:
        ax.text(b.get_x() + b.get_width() / 2, max(ys) * 0.03, "belum\nterbit",
                ha="center", va="bottom", fontsize=8, color=HEX(RUST), fontweight="bold")
    ax.set_ylabel("hektar", fontsize=10)
    ax.grid(axis="y", color=HEX(LINE), linewidth=0.7)
    ax.set_axisbelow(True)
    ax.margins(y=0.22)
    return _finish(fig, name)


# ---------------------------------------------------------------- derivasi
def conf_buckets(rows):
    lo = nom = hi = 0
    for r in rows:
        c = str(r["confidence"]).strip().lower()
        n = r["n"]
        if c == "h":
            hi += n
        elif c == "n":
            nom += n
        elif c in ("l", "(kosong)", ""):
            lo += n
        elif c.isdigit():
            v = int(c)
            hi += n if v >= 80 else 0
            nom += n if 30 <= v < 80 else 0
            lo += n if v < 30 else 0
        else:
            nom += n
    return lo, nom, hi


KEL_ORDER = ["Konservasi", "Lindung", "Produksi", "Non-Kawasan Hutan"]
KEL_COL = {"Konservasi": GREEN, "Lindung": SLATE, "Produksi": EMBER,
           "Non-Kawasan Hutan": GOLD}


def kel_series(rows, key="n"):
    d = {r["kelompok"]: float(r[key]) for r in rows}
    labels, vals, cols = [], [], []
    for k in KEL_ORDER:
        if k in d:
            labels.append(k)
            vals.append(d[k])
            cols.append(HEX(KEL_COL[k]))
    for r in rows:
        if r["kelompok"] not in KEL_ORDER and float(r[key]) > 0:
            labels.append(r["kelompok"])
            vals.append(float(r[key]))
            cols.append(HEX(LINE))
    return labels, vals, cols


def chart_kel(rows, name, *, key="n", xlabel="", fmt=ribu, figsize=(7.2, 2.6)):
    labels, vals, cols = kel_series(rows, key)
    fig, ax = plt.subplots(figsize=figsize)
    y = list(range(len(labels)))
    ax.barh(y, vals, color=cols, height=0.6)
    for i, v in enumerate(vals):
        ax.text(v, i, " " + fmt(v), va="center", fontsize=9, color=HEX(INK),
                fontweight="bold")
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=10)
    ax.invert_yaxis()
    ax.set_xlabel(xlabel, fontsize=9.5)
    ax.grid(axis="x", color=HEX(LINE), linewidth=0.7)
    ax.set_axisbelow(True)
    ax.margins(x=0.2)
    return _finish(fig, name)


# ---------------------------------------------------------------- build
def build(data_path: str, out_path: str):
    D = json.loads(Path(data_path).read_text(encoding="utf-8"))
    P, H, K, KL, EX = D["profil"], D["hotspot"], D["kompleks"], D["klhk"], D["extra"]
    per = D["periode"]
    per_end_id = tanggal_id(per["end"])
    per_start_id = tanggal_id(per["start"])
    ps_n = next((r["n"] for r in P["kps_by_layer"] if r["layer_key"] == "psagustus2026"), 0)
    ha_n = next((r["n"] for r in P["kps_by_layer"] if r["layer_key"] == "HUTAN_ADAT_APR26"), 0)
    mvd = {r["m"]: r["n"] for r in H["hs_monthly"]}
    lo, nom, hi = conf_buckets(H["hs_by_confidence"])
    tot_c = lo + nom + hi
    viirs = sum(r["n"] for r in H["hs_by_source"] if r["source"].startswith("VIIRS"))
    kl_last_m = KL["total"]["bln_akhir"].split("-")[1] if KL["total"].get("bln_akhir") else "07"
    kl_last_label = BULAN_PANJANG[int(kl_last_m)]

    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H
    pg = [0]

    def nextpg():
        pg[0] += 1
        return pg[0]

    # ============================================================ SAMPUL
    s = blank(prs)
    _bg(s, INK)
    _rect(s, Inches(0), Inches(0), EMU_W, Inches(0.16), EMBER)
    _rect(s, Inches(0.75), Inches(2.3), Inches(0.9), Inches(0.09), EMBER)
    _txt(s, Inches(0.75), Inches(1.62), Inches(11.6), Inches(0.4),
         [[("ETA SENEU — SISTEM PERINGATAN DINI KEBAKARAN HUTAN DAN LAHAN", 12.5, True, EMBER)]])
    _txt(s, Inches(0.72), Inches(2.55), Inches(12.1), Inches(2.0),
         [[("Analisis Kebakaran Hutan dan Lahan", 38, True, WHITE)],
          [("pada Areal Perhutanan Sosial dan Hutan Adat", 22, True,
            RGBColor(0xD8, 0xDD, 0xE2))]], line_spacing=1.08)
    _txt(s, Inches(0.75), Inches(4.55), Inches(12), Inches(0.5),
         [[("Rekapitulasi Titik Panas  ·  Kompleks Kebakaran  ·  Luas Areal "
            "Kebakaran menurut Data Kementerian Kehutanan", 13, True,
            RGBColor(0xB8, 0xC0, 0xC8))]])
    _txt(s, Inches(0.75), Inches(5.25), Inches(12), Inches(0.45),
         [[(f"Periode {per_start_id} s.d. {per_end_id}", 13, False,
            RGBColor(0xB8, 0xC0, 0xC8))]])
    _txt(s, Inches(0.75), Inches(5.95), Inches(11.9), Inches(1.0),
         [[("Cakupan: " + ribu(P["kps_total"]) + " unit KPS dan Hutan Adat aktif pada "
            + str(P["n_prov"]) + " provinsi (" + ribu(ps_n) + " Perhutanan Sosial + "
            + ribu(ha_n) + " Hutan Adat).", 10.5, False, RGBColor(0x9A, 0xA4, 0xAE))],
          [("Sumber: titik panas NASA FIRMS (VIIRS 375 m + MODIS 1 km); rekapitulasi "
            "resmi Kementerian Kehutanan “Areal Kebakaran Hutan dan Lahan”. "
            "Data diolah pada " + per_end_id + ".", 10.5, False, RGBColor(0x9A, 0xA4, 0xAE))]],
         line_spacing=1.2)
    footer(s, nextpg(), dark=True)

    # ============================================================ RINGKASAN
    s = content(prs, "Ringkasan Eksekutif",
                "Titik panas terkonsentrasi pada Agustus 2026; angka luas resmi "
                "baru tersedia s.d. " + kl_last_label, nextpg())
    stat_cards(s, [
        (ribu(H["hs_total"]), "titik panas terpantau\nsepanjang 2026 (s.d. " +
         per_end_id + ")", GOLD),
        (ribu(mvd.get(8, 0)), "di antaranya terjadi\npada Agustus 2026 (" +
         pct(mvd.get(8, 0), H["hs_total"]) + ")", RUST),
        (ribu(K["n_kompleks"]), "kompleks kebakaran\ndalam 30 hari terakhir", EMBER),
        (ribu(H["kps_with_hotspot"]), "unit KPS/Hutan Adat\nterpapar titik panas", SLATE),
        (ha(KL["total"]["ha"]) + " ha", "luas areal kebakaran resmi\n(Jan–" +
         kl_last_label[:3] + " · " + ribu(KL["total"]["kps"]) + " unit)", INK_SOFT),
    ])
    bullets(s, [
        ("Puncak musiman. ", "Sebanyak " + pct(mvd.get(8, 0), H["hs_total"]) +
         " titik panas tahun 2026 terjadi pada bulan Agustus (" + ribu(mvd.get(8, 0)) +
         " dari " + ribu(H["hs_total"]) + " titik), meningkat sekitar " +
         f"{mvd.get(8,0)/max(mvd.get(7,1),1):.1f}".replace(".", ",") +
         " kali lipat dibandingkan Juli."),
        ("Terkonsentrasi wilayah. ", "Kalimantan Barat, Papua Selatan, dan Kalimantan "
         "Tengah menyumbang bagian terbesar titik panas nasional pada areal KPS/Hutan Adat."),
        ("Kelas keyakinan didominasi tingkat sedang. ", pct(nom, tot_c) +
         " titik berkelas nominal/sedang dan hanya " + pct(hi, tot_c) +
         " berkeyakinan tinggi — mayoritas memerlukan verifikasi lanjutan."),
        ("Data luas resmi tertinggal musim. ", "Rekapitulasi Kementerian Kehutanan baru "
         "sampai " + kl_last_label + " 2026 (" + ha(KL["total"]["ha"]) +
         " ha); bulan Agustus–September belum diterbitkan."),
        ("Hotspot bukan ukuran luas. ", "Titik panas berfungsi sebagai penanda dini; "
         "penilaian luas dan dampak tetap mengacu pada data resmi dan verifikasi lapangan."),
    ], y=Inches(3.75), size=10.5)

    # ============================================================ METODOLOGI
    s = content(prs, "Sumber Data dan Metodologi",
                "Dua sumber data dengan peran yang berbeda", nextpg())
    cx, cw = Inches(0.78), Inches(3.82)
    cards = [
        ("Titik panas — NASA FIRMS", EMBER,
         "Sensor VIIRS (S-NPP, NOAA-20, NOAA-21) resolusi 375 m dan MODIS 1 km, "
         "disinkronkan berkala. Menandai anomali suhu permukaan sebagai indikasi dini, "
         "bukan ukuran luas. Titik dipetakan ke poligon KPS/Hutan Adat melalui "
         "analisis spasial (uji titik-dalam-poligon)."),
        ("Kompleks kebakaran — turunan", SLATE,
         "Pengelompokan titik panas yang berdekatan secara ruang (± 2 km) dan waktu "
         "(± 48 jam) menjadi satu “kompleks” (metode ST-DBSCAN, preset sensitivitas "
         "sedang). Menu “Kompleks Kebakaran” pada sistem ETA SENEU. Jendela analisis: "
         "30 hari terakhir (" + tanggal_id(K["window_start"]) + " s.d. " +
         tanggal_id(K["window_end"]) + ")."),
        ("Luas areal kebakaran — Kementerian Kehutanan", GREEN,
         "Dokumen resmi “Areal Kebakaran Hutan dan Lahan”, akurasi tinggi/sedang. "
         "Terbit tidak terjadwal. Untuk 2026 baru tersedia Januari–" + kl_last_label +
         ". Angka luas tidak dijumlahkan dengan jumlah titik panas — satuan dan makna "
         "berbeda."),
    ]
    for i, (title_c, acc, body) in enumerate(cards):
        x = Emu(int(cx + i * (cw + Inches(0.2))))
        _rect(s, x, Inches(1.95), cw, Inches(3.7), WHITE, line=LINE, line_w=Pt(1), rounded=True)
        _rect(s, x, Inches(1.95), cw, Inches(0.09), acc)
        _txt(s, Emu(int(x + Inches(0.26))), Inches(2.2), Emu(int(cw - Inches(0.52))),
             Inches(0.8), [[(title_c, 12.5, True, INK)]], line_spacing=1.05)
        _txt(s, Emu(int(x + Inches(0.26))), Inches(3.15), Emu(int(cw - Inches(0.52))),
             Inches(2.4), [[(body, 9.5, False, INK_SOFT)]], line_spacing=1.2)
    bullets(s, [
        "Seluruh angka merupakan irisan dengan poligon KPS dan Hutan Adat aktif "
        "(psagustus2026 + HUTAN_ADAT_APR26) — bukan total karhutla wilayah.",
        "Atribusi fungsi kawasan hutan mengacu pada peta KWSHUTAN_AR_250K "
        "Kementerian Kehutanan (kelompok: Konservasi / Lindung / Produksi / Non-Kawasan Hutan).",
        "Data ditarik READ-ONLY dari basis data operasional ETA SENEU pada " + per_end_id + ".",
    ], y=Inches(5.85), size=9.5, gap=Pt(4))

    # ============================================================ BAGIAN I
    divider(prs, "Bagian I", "Rekapitulasi Titik Panas (Hotspot)",
            "Periode 1 Januari s.d. " + per_end_id + "  ·  menurut provinsi dan "
            "wilayah kerja Balai Perhutanan Sosial  ·  " + ribu(H["hs_total"]) +
            " titik panas terpantau", nextpg())

    # ---- I.1 tren bulanan
    s = content(prs, "Bagian I · Tren Bulanan",
                "Aktivitas titik panas rendah sepanjang semester I, melonjak pada "
                "triwulan III", nextpg())
    place(s, chart_monthly([(r["m"], r["n"]) for r in H["hs_monthly"]],
                           "hs_monthly.png", ylabel="titik panas"),
          Inches(0.72), Inches(1.85), Inches(11.9))
    bullets(s, [
        (ribu(H["hs_total"]) + " titik panas sepanjang 2026: ",
         "Januari–Juni di bawah " + ribu(mvd.get(4, 0)) + " titik/bulan, Juli " +
         ribu(mvd.get(7, 0)) + " titik, Agustus " + ribu(mvd.get(8, 0)) +
         " titik, dan September (berjalan) " + ribu(mvd.get(9, 0)) + " titik."),
        ("Anomali April (" + ribu(mvd.get(4, 0)) + " titik): ",
         "kenaikan singkat pada awal kemarau, terkait pembersihan lahan di beberapa sentra."),
        ("Cakupan keterpaparan: ", ribu(H["kps_with_hotspot"]) + " unit KPS/Hutan Adat "
         "terpapar (" + pct(H["kps_with_hotspot"], P["kps_total"]) + " dari " +
         ribu(P["kps_total"]) + " unit aktif)."),
    ], y=Inches(5.3), size=10.5)

    # ---- I.2 per provinsi
    s = content(prs, "Bagian I · Sebaran Menurut Provinsi",
                "Empat provinsi menanggung lebih dari 80 persen titik panas nasional",
                nextpg())
    pv = H["hs_by_prov"][:15]
    place(s, chart_hbar([r["nama_prov"] for r in pv], [r["n"] for r in pv],
                        "hs_prov.png", color=EMBER, xlabel="titik panas (2026)",
                        figsize=(7.6, 5.0)),
          Inches(0.72), Inches(1.85), Inches(7.6))
    top4 = sum(r["n"] for r in H["hs_by_prov"][:4])
    _txt(s, Inches(8.5), Inches(1.95), Inches(4.2), Inches(4.9), [
        [("Catatan", 12.5, True, INK)],
        [("Kalimantan Barat (" + ribu(pv[0]["n"]) + "), Papua Selatan (" +
          ribu(pv[1]["n"]) + "), Kalimantan Tengah (" + ribu(pv[2]["n"]) +
          "), dan Papua Barat (" + ribu(pv[3]["n"]) + ") ", 10, True, INK),
         ("mencakup " + pct(top4, H["hs_total"]) + " titik panas nasional pada areal "
          "KPS/Hutan Adat.", 10, False, INK_SOFT)],
        [("Provinsi lain di bawah " + ribu(pv[4]["n"]) + " titik.", 10, False, INK_SOFT)],
        [("Titik panas menyebar pada " + str(len(H["hs_by_prov"])) +
          " provinsi; " + ribu(H["hs_in_polygon"]) + " dari " + ribu(H["hs_total"]) +
          " titik berada di dalam poligon KPS/Hutan Adat.", 10, False, INK_SOFT)],
    ], line_spacing=1.22, space_after=Pt(9))

    # ---- I.3 per Balai PS
    s = content(prs, "Bagian I · Sebaran Menurut Wilayah Kerja Balai",
                "Wilayah kerja Balai PS Banjarbaru dan Manokwari menjadi beban utama",
                nextpg())
    wv = [r for r in H["hs_by_wilker"] if r["wilker_bps"] != "(tidak tercatat)"][:12]
    place(s, chart_hbar([r["wilker_bps"].replace("Balai PS ", "") for r in wv],
                        [r["n"] for r in wv], "hs_wilker.png", color=SLATE,
                        xlabel="titik panas (2026)", figsize=(7.7, 4.8)),
          Inches(0.72), Inches(1.85), Inches(7.7))
    _txt(s, Inches(8.6), Inches(1.95), Inches(4.1), Inches(4.9), [
        [("Pembacaan", 12.5, True, INK)],
        [("Kolom wilker_bps pada metadata poligon merujuk pada Balai/wilayah kerja "
          "Perhutanan Sosial. Satu wilayah kerja dapat mencakup beberapa provinsi "
          "(mis. Banjarbaru mencakup Kalimantan Barat, Tengah, dan Selatan).",
          9.5, False, INK_SOFT)],
        [("Sebagian poligon belum tercatat wilayah kerjanya; angka pada grafik ini "
          "hanya untuk poligon yang tercatat.", 9.5, False, INK_SOFT)],
        [("Rincian lengkap provinsi × wilayah kerja disajikan pada halaman berikut.",
          9.5, False, INK_SOFT)],
    ], line_spacing=1.2, space_after=Pt(8))

    # ---- I.4 rincian provinsi x Balai (DAFTAR LENGKAP, dipaginasi)
    pw_all = sorted(H["hs_by_prov_wilker"], key=lambda r: -r["n"])
    per_page = 13
    n_page = (len(pw_all) + per_page - 1) // per_page
    for pi in range(n_page):
        chunk = pw_all[pi * per_page:(pi + 1) * per_page]
        rk = pi * per_page
        s = content(prs, "Bagian I · Rincian Provinsi dan Wilayah Kerja Balai",
                    "Titik panas menurut provinsi dan wilayah kerja Balai PS — daftar "
                    "lengkap (" + str(pi + 1) + "/" + str(n_page) + ")", nextpg())
        table(s, ["#", "Provinsi", "Wilayah Kerja Balai PS", "Titik Panas"],
              [[str(rk + j + 1), r["nama_prov"], r["wilker_bps"], ribu(r["n"])]
               for j, r in enumerate(chunk)],
              x=Inches(0.78), y=Inches(1.9), w=Inches(11.9), row_h=Inches(0.34), size=10.5,
              col_w=[Inches(0.7), Inches(4.2), Inches(4.6), Inches(2.4)])
        if pi == n_page - 1:
            _txt(s, Inches(0.78), Inches(6.6), Inches(11.9), Inches(0.4),
                 [[("Total " + str(len(pw_all)) + " kombinasi provinsi × wilayah kerja. "
                    "Nilai “(tidak tercatat)” = poligon tanpa isian wilker_bps.",
                    9, False, INK_SOFT)]])

    # ---- I.5 kualitas deteksi
    s = content(prs, "Bagian I · Kualitas Deteksi",
                "Mayoritas titik panas berkelas keyakinan sedang", nextpg())
    place(s, chart_hbar(["Sedang / nominal", "Rendah", "Tinggi"], [nom, lo, hi],
                        "hs_conf.png", color=GOLD, xlabel="titik panas (2026)",
                        figsize=(7.0, 2.3)),
          Inches(0.72), Inches(1.9), Inches(7.0))
    src = H["hs_by_source"]
    place(s, chart_hbar([r["source"] for r in src], [r["n"] for r in src],
                        "hs_src.png", color=SLATE, xlabel="titik panas (2026)",
                        figsize=(7.0, 2.6)),
          Inches(0.72), Inches(4.35), Inches(7.0))
    _txt(s, Inches(8.3), Inches(1.95), Inches(4.4), Inches(4.9), [
        [("Interpretasi", 12.5, True, INK)],
        [(pct(nom, tot_c) + " titik berkelas sedang dan " + pct(lo, tot_c) +
          " rendah; hanya " + pct(hi, tot_c) + " (± " + ribu(hi) + " titik) berkeyakinan tinggi.",
          10, False, INK_SOFT)],
        [("Titik berkelas sedang wajar pada kebakaran kecil/permukaan, namun juga dapat "
          "berasal dari lahan panas atau sumber termal lain — memerlukan konfirmasi "
          "citra/lapangan.", 10, False, INK_SOFT)],
        [("Sensor VIIRS menyumbang " + pct(viirs, H["hs_total"]) + " deteksi; sisanya MODIS.",
          10, False, INK_SOFT)],
    ], line_spacing=1.2, space_after=Pt(8))

    # ---- I.6 fungsi kawasan hutan
    s = content(prs, "Bagian I · Menurut Fungsi Kawasan Hutan",
                "Titik panas dominan pada kawasan Hutan Produksi", nextpg())
    place(s, chart_kel(H["hs_by_kawasan"], "hs_kel.png", key="n",
                       xlabel="titik panas (2026)", figsize=(8.0, 2.8)),
          Inches(0.72), Inches(2.0), Inches(8.0))
    kd = {r["kelompok"]: r["n"] for r in H["hs_by_kawasan"]}
    tot_kel = sum(kd.values())
    _txt(s, Inches(9.0), Inches(2.0), Inches(3.7), Inches(4.5), [
        [("Catatan", 12.5, True, INK)],
        [("Hutan Produksi " + pct(kd.get("Produksi", 0), tot_kel) + ", Hutan Lindung " +
          pct(kd.get("Lindung", 0), tot_kel) + ".", 10, False, INK_SOFT)],
        [("Titik pada Hutan Lindung dan Konservasi menjadi prioritas penanganan karena "
          "nilai perlindungan kawasannya lebih tinggi.", 10, False, INK_SOFT)],
        [("Atribusi mengacu peta KWSHUTAN_AR_250K Kementerian Kehutanan.",
          9.5, False, INK_SOFT)],
    ], line_spacing=1.2, space_after=Pt(8))
    bullets(s, [
        "Kelompok “Non-Kawasan Hutan” mencakup Areal Penggunaan Lain (APL) di dalam "
        "batas poligon KPS/Hutan Adat.",
    ], y=Inches(5.35), size=10)

    # ---- I.7 pola waktu
    s = content(prs, "Bagian I · Pola Waktu Kejadian",
                "Deteksi memuncak pada siang hari serta pada lintasan tengah malam",
                nextpg())
    place(s, chart_hours([(r["h"], r["n"]) for r in H["hs_by_hour"]], "hs_hours.png"),
          Inches(0.72), Inches(1.9), Inches(11.9))
    bullets(s, [
        ("Puncak siang (pukul 11.00–14.00 WIB): ", "sesuai pola pembakaran lahan dan "
         "kondisi permukaan terpanas — jendela kewaspadaan operasional tertinggi."),
        ("Puncak dini hari (pukul 00.00–01.00 WIB): ", "lintasan malam sensor VIIRS; "
         "menandakan bara/kebakaran yang masih menyala setelah gelap."),
        ("Implikasi: ", "patroli dan kesiapsiagaan pemadaman difokuskan pada rentang "
         "siang, dengan pemantauan lanjutan pada malam hari."),
    ], y=Inches(4.95), size=10.5)

    # ============================================================ BAGIAN II
    divider(prs, "Bagian II", "Kompleks Kebakaran — 30 Hari Terakhir",
            tanggal_id(K["window_start"]) + " s.d. " + tanggal_id(K["window_end"]) +
            "  ·  metode ST-DBSCAN (sensitivitas sedang)  ·  " + ribu(K["n_kompleks"]) +
            " kompleks kebakaran teridentifikasi", nextpg())

    # ---- II.1 ringkasan kompleks
    s = content(prs, "Bagian II · Ringkasan",
                "Kompleks kebakaran memusatkan ribuan titik lepas menjadi satuan kejadian",
                nextpg())
    stat_cards(s, [
        (ribu(K["n_kompleks"]), "kompleks kebakaran\ndalam 30 hari terakhir", EMBER),
        (ribu(K["clustered_hotspots"]), "titik panas tergabung\ndalam kompleks", GOLD),
        (ribu(K["kompleks_besar"]), "kompleks berukuran besar\n(≥ 10 titik)", RUST),
        (ribu(K["kps_terlibat"]), "unit KPS/Hutan Adat\nterlibat kompleks", SLATE),
        (f"{K['durasi_jam_rata']:.0f}".replace(".", ",") + " jam",
         "durasi rata-rata kompleks\n(titik pertama–terakhir)", INK_SOFT),
    ])
    bullets(s, [
        ("Definisi. ", "Satu “kompleks kebakaran” adalah gugusan titik panas yang "
         "berdekatan ruang (± 2 km) dan waktu (± 48 jam) — mendekati satu kejadian "
         "kebakaran nyata, bukan sekadar jumlah titik satelit."),
        ("Dari " + ribu(K["total_hotspots_in_range"]) + " titik pada jendela 30 hari, " +
         ribu(K["clustered_hotspots"]) + " titik ", "membentuk kompleks; sisanya "
         "terdeteksi menyendiri (belum memenuhi ambang kepadatan)."),
        ("Durasi. ", "Rata-rata " + f"{K['durasi_jam_rata']:.0f}".replace(".", ",") +
         " jam; kompleks terpanjang mencapai " +
         f"{K['durasi_jam_maks']/24:.0f}".replace(".", ",") + " hari — indikasi "
         "kebakaran yang berlangsung menerus dan perlu penanganan segera."),
    ], y=Inches(3.7), size=10.5)

    # ---- II.2 KPS dengan kompleks terbanyak
    s = content(prs, "Bagian II · KPS dengan Kompleks Kebakaran Terbanyak",
                "Sepuluh KPS/Hutan Adat dengan kompleks kebakaran terbanyak (30 hari "
                "terakhir)", nextpg())
    tk = K["top_kps"][:10]
    table(s, ["Lembaga / KPS", "Provinsi", "Kab./Kota", "Skema", "Kompleks", "Titik"],
          [[(r["lembaga"] or "—")[:34], r["nama_prov"] or "—", r["nama_kab"] or "—",
            r["skema"] or "—", ribu(r["kompleks"]), ribu(r["titik"])] for r in tk],
          x=Inches(0.78), y=Inches(1.9), w=Inches(11.9), row_h=Inches(0.36), size=10,
          col_w=[Inches(3.9), Inches(2.0), Inches(1.9), Inches(1.35), Inches(1.4),
                 Inches(1.35)])
    bullets(s, [
        "Kolom “Titik” = jumlah titik panas dalam kompleks yang berpusat pada KPS "
        "bersangkutan selama 30 hari terakhir.",
        "Klaster terkonsentrasi di Kalimantan Barat (Kabupaten Ketapang) dan "
        "Papua Selatan (Kabupaten Mappi) — menjadi prioritas verifikasi lapangan.",
    ], y=Inches(6.0), size=9.5, gap=Pt(4))

    # ============================================================ BAGIAN III
    divider(prs, "Bagian III", "Luas Areal Kebakaran menurut Data Kementerian Kehutanan",
            "Rekapitulasi resmi “Areal Kebakaran Hutan dan Lahan”  ·  Januari–" +
            kl_last_label + " 2026  ·  " + ha(KL["total"]["ha"]) + " ha pada " +
            ribu(KL["total"]["kps"]) + " unit KPS/Hutan Adat", nextpg())

    # ---- III.1 bulanan
    s = content(prs, "Bagian III · Luas Areal Kebakaran Bulanan",
                "Rekapitulasi resmi baru sampai " + kl_last_label +
                "; puncak Agustus belum tercatat", nextpg())
    kmv = [(r["month"], float(r["ha"])) for r in KL["monthly"]]
    kmvd = dict(kmv)
    place(s, chart_klhk_monthly(kmv, "klhk_monthly.png"),
          Inches(0.72), Inches(1.85), Inches(11.9))
    hi_m = max(kmv, key=lambda x: x[1])
    bullets(s, [
        (ha(KL["total"]["ha"]) + " ha pada " + ribu(KL["total"]["kps"]) +
         " unit KPS/Hutan Adat (Jan–" + kl_last_label[:3] + " 2026). ",
         "Luas terbesar pada " + BULAN_PANJANG[hi_m[0]] + " (" + ha(hi_m[1]) + " ha)."),
        ("Keterbatasan periode. ", "Bulan Agustus 2026 — yang menunjukkan lonjakan "
         "titik panas tertinggi — belum tercakup; angka luas final akan menyusul "
         "setelah rekapitulasi Kementerian Kehutanan berikutnya terbit."),
        ("Sampai saat itu, ", "pengukuran dampak Agustus–September bersifat indikatif "
         "berbasis titik panas dan kompleks kebakaran."),
    ], y=Inches(5.25), size=10.5)

    # ---- III.2 per provinsi
    s = content(prs, "Bagian III · Menurut Provinsi",
                "Riau, Kalimantan Barat, dan Maluku menanggung luas terbakar terbesar",
                nextpg())
    kp = KL["by_prov"][:10]
    place(s, chart_hbar([r["nama_prov"] for r in kp], [float(r["ha"]) for r in kp],
                        "klhk_prov.png", color=SLATE, xlabel="hektar (Jan–" +
                        kl_last_label[:3] + " 2026)", fmt=ha, figsize=(7.5, 4.7)),
          Inches(0.72), Inches(1.9), Inches(7.5))
    table(s, ["Provinsi", "Unit", "Hektar"],
          [[r["nama_prov"], ribu(r["kps"]), ha(r["ha"])] for r in kp[:10]],
          x=Inches(8.5), y=Inches(1.95), w=Inches(4.2), row_h=Inches(0.44), size=9.5,
          col_w=[Inches(2.1), Inches(0.9), Inches(1.2)])

    # ---- III.3 per fungsi kawasan
    s = content(prs, "Bagian III · Menurut Fungsi Kawasan Hutan",
                "Sebagian besar luas terbakar resmi berada pada kawasan Hutan Produksi",
                nextpg())
    if KL.get("by_kawasan"):
        place(s, chart_kel(KL["by_kawasan"], "klhk_kel.png", key="ha",
                           xlabel="hektar (Jan–" + kl_last_label[:3] + " 2026)", fmt=ha,
                           figsize=(8.0, 2.8)),
              Inches(0.72), Inches(2.0), Inches(8.0))
        kkd = {r["kelompok"]: float(r["ha"]) for r in KL["by_kawasan"]}
        tot_kk = sum(kkd.values())
        _txt(s, Inches(9.0), Inches(2.0), Inches(3.7), Inches(4.3), [
            [("Catatan", 12.5, True, INK)],
            [("Hutan Produksi " + pct(kkd.get("Produksi", 0), tot_kk) +
              ", Hutan Lindung " + pct(kkd.get("Lindung", 0), tot_kk) + " dari total " +
              ha(tot_kk) + " ha.", 10, False, INK_SOFT)],
            [("Diperoleh dengan mengiris geometri areal terbakar resmi terhadap peta "
              "fungsi kawasan hutan Kementerian Kehutanan.", 9.5, False, INK_SOFT)],
        ], line_spacing=1.2, space_after=Pt(8))
    else:
        _txt(s, Inches(0.78), Inches(2.2), Inches(11.9), Inches(1.0),
             [[("Rincian per fungsi kawasan tidak tersedia pada data saat ini.",
                11, False, INK_SOFT)]])

    # ============================================================ BAGIAN IV
    divider(prs, "Bagian IV", "Analisis Lanjutan",
            "Kepadatan titik panas per satuan luas  ·  KPS dengan kebakaran berulang  ·  "
            "arah tindak lanjut", nextpg())

    # ---- IV.1 kepadatan per 1000 ha
    s = content(prs, "Bagian IV · Kepadatan Titik Panas per 1.000 Ha",
                "Dinormalkan terhadap luas, sejumlah unit kecil menunjukkan intensitas "
                "tertinggi", nextpg())
    dt = EX["density_top"][:10]
    table(s, ["Lembaga / KPS", "Provinsi", "Luas (ha)", "Titik", "Titik / 1.000 ha"],
          [[(r["lembaga"] or "—")[:34], r["nama_prov"] or "—", ha(r["luas_ha"]),
            ribu(r["n"]), ha(r["per_1000ha"])] for r in dt],
          x=Inches(0.78), y=Inches(1.9), w=Inches(11.9), row_h=Inches(0.36), size=10,
          col_w=[Inches(4.1), Inches(2.3), Inches(1.8), Inches(1.5), Inches(2.2)])
    bullets(s, [
        "Hanya poligon berluas ≥ 50 ha dan ≥ 5 titik panas yang diperhitungkan, "
        "agar rasio tidak bias pada poligon sangat kecil.",
        "Daftar ini melengkapi peringkat jumlah absolut: unit berluas kecil dengan "
        "kepadatan tinggi berisiko terbakar hampir menyeluruh.",
    ], y=Inches(5.95), size=9.5, gap=Pt(4))

    # ---- IV.2 kebakaran berulang
    s = content(prs, "Bagian IV · KPS dengan Kebakaran Berulang",
                str(EX["chronic_count"]) + " KPS/Hutan Adat terpapar titik panas pada "
                "≥ 2 bulan berbeda sepanjang 2026", nextpg())
    ch = EX["chronic"][:10]
    table(s, ["Lembaga / KPS", "Provinsi", "Bln Terpapar", "Titik", "Rentang"],
          [[(r["lembaga"] or "—")[:32], r["nama_prov"] or "—", ribu(r["bulan_kena"]),
            ribu(r["n"]), r["d0"][5:] + " s.d. " + r["d1"][5:]] for r in ch],
          x=Inches(0.78), y=Inches(1.9), w=Inches(11.9), row_h=Inches(0.36), size=9.5,
          col_w=[Inches(3.9), Inches(2.3), Inches(1.7), Inches(1.3), Inches(2.7)])
    bullets(s, [
        "Unit dengan keterpaparan ≥ 6 bulan menandakan tekanan kebakaran yang menahun — "
        "sasaran pembinaan dan pendampingan jangka panjang, bukan respons sesaat.",
    ], y=Inches(5.95), size=9.5, gap=Pt(4))

    # ---- IV.3 rekomendasi
    s = content(prs, "Bagian IV · Rekomendasi dan Tindak Lanjut",
                "Prioritas penanganan dan pendampingan", nextpg())
    _txt(s, Inches(0.78), Inches(1.8), Inches(6.0), Inches(0.32),
         [[("ARAH TINDAK LANJUT", 11.5, True, EMBER)]])
    bullets(s, [
        "Menetapkan Agustus–September sebagai periode siaga tertinggi untuk KPS di "
        "Kalimantan Barat, Papua Selatan, Kalimantan Tengah, Papua Barat, dan Riau.",
        "Mengarahkan verifikasi lapangan ke KPS dengan kompleks kebakaran terbanyak "
        "(Bagian II) dan kebakaran berulang (Bagian IV).",
        "Mengintegrasikan rekapitulasi luas Kementerian Kehutanan begitu diterbitkan; "
        "sampai saat itu dampak dilaporkan indikatif berbasis titik panas.",
        "Tidak menyimpulkan luas kebakaran dari jumlah titik panas; menyandingkan "
        "dengan citra resolusi tinggi dan laporan pendamping KPS.",
        "Memprioritaskan penanganan titik panas pada kawasan Hutan Lindung dan Konservasi.",
    ], x=Inches(0.78), y=Inches(2.2), w=Inches(6.0), h=Inches(4.6), size=9.5, gap=Pt(7))
    _txt(s, Inches(7.1), Inches(1.8), Inches(5.6), Inches(0.32),
         [[("PENDAMPING KPS PRIORITAS", 11.5, True, EMBER)]])
    pp = D.get("pendamping_prioritas", [])[:9]
    if pp:
        table(s, ["KPS", "Provinsi", "Pendamping"],
              [[(r["lembaga"] or "—")[:22], (r["nama_prov"] or "—")[:18],
                (r["pendamping"] or "—")[:24]] for r in pp],
              x=Inches(7.1), y=Inches(2.2), w=Inches(5.65), row_h=Inches(0.42), size=8.5,
              col_w=[Inches(1.9), Inches(1.6), Inches(2.15)])
        _txt(s, Inches(7.1), Inches(6.5), Inches(5.65), Inches(0.5),
             [[("Sumber: Master Data Pendamping PS 2026. Tanda “—” = nama pendamping "
                "belum terpetakan (perbedaan penulisan nama lembaga).",
                8, False, INK_SOFT)]], line_spacing=1.15)
    else:
        _txt(s, Inches(7.1), Inches(2.2), Inches(5.6), Inches(1.0),
             [[("Data pendamping tidak tersedia.", 10, False, INK_SOFT)]])

    # ============================================================ CATATAN
    s = content(prs, "Penutup", "Catatan, Batasan, dan Sumber Data", nextpg())
    bullets(s, [
        "Cakupan terbatas pada irisan dengan poligon KPS dan Hutan Adat aktif "
        "(psagustus2026 + HUTAN_ADAT_APR26) — bukan total kebakaran hutan dan lahan wilayah.",
        "Titik panas adalah anomali termal NASA FIRMS (resolusi 375 m VIIRS / 1 km "
        "MODIS); mayoritas berkelas keyakinan sedang. Bukan konfirmasi kebakaran dan "
        "bukan ukuran luas.",
        "Pemetaan titik panas ke poligon menggunakan analisis spasial saat kueri "
        "(uji titik-dalam-poligon); tabel relasi bawaan tidak dipakai karena "
        "pemutakhirannya tertinggal.",
        "“Kompleks kebakaran” dihitung dengan metode ST-DBSCAN preset sensitivitas "
        "sedang (radius ± 2 km, selisih waktu ± 48 jam) pada jendela 30 hari terakhir.",
        "Luas areal kebakaran bersumber dari rekapitulasi resmi Kementerian Kehutanan "
        "“Areal Kebakaran Hutan dan Lahan” akurasi tinggi/sedang; data 2026 baru "
        "Januari–" + kl_last_label + ".",
        "Fungsi kawasan hutan mengacu peta KWSHUTAN_AR_250K Kementerian Kehutanan.",
        "Seluruh data ditarik READ-ONLY dari basis data operasional ETA SENEU pada " +
        per_end_id + ".",
    ], y=Inches(1.95), size=10, gap=Pt(8))

    prs.save(out_path)
    print("tersimpan:", out_path, "|", pg[0], "slide")


if __name__ == "__main__":
    build(sys.argv[1], sys.argv[2])
